# ============================================================================
# VetMis Technology Slide Generator
# Создаёт слайд для PowerPoint с информацией о технологичности проекта
# ============================================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# ЦВЕТОВАЯ ПАЛИТРА
# ============================================================================
COLORS = {
    'primary_blue': RGBColor(46, 134, 171),    # #2E86AB
    'accent_pink': RGBColor(162, 59, 114),     # #A23B72
    'bg_light': RGBColor(248, 249, 250),       # #F8F9FA
    'text_dark': RGBColor(45, 55, 72),         # #2D3748
    'white': RGBColor(255, 255, 255),
    'green': RGBColor(72, 187, 120),           # #48BB78
}

# ============================================================================
# КОНТЕНТ СЛАЙДА
# ============================================================================
SLIDE_CONTENT = {
    'title': 'Технологичность стартап-проекта VetMis',
    
    'column1': {
        'icon': '🛠️',
        'title': 'Технологический стек',
        'items': [
            'Backend: Python 3.10+ / Django 5.0+',
            'Frontend: HTML5/CSS3 (vanilla)',
            'База данных: SQLite → PostgreSQL',
            'API: Django REST Framework',
            'Библиотеки: qrcode, Pillow'
        ]
    },
    
    'column2': {
        'icon': '🏗️',
        'title': 'Архитектурные преимущества',
        'items': [
            '✓ Модульная MVT-архитектура',
            '✓ Ролевая модель (3 типа пользователей)',
            '✓ API-first подход для интеграций',
            '✓ Генерация QR-паспортов',
            '✓ Загрузка и обработка документов'
        ]
    },
    
    'column3': {
        'icon': '📈',
        'title': 'Масштабируемость',
        'items': [
            '→ PostgreSQL + Redis + Celery',
            '→ Микросервисная архитектура',
            '→ Интеграция с ЕГИС и «Меркурий»',
            '→ Мобильные приложения (RN/Flutter)',
            '→ Платёжные системы и HTTPS'
        ]
    },
    
    'footer': '💡 Уникальное преимущество: MVP с архитектурой, готовой к промышленной эксплуатации — баланс между скоростью разработки и технологической зрелостью.'
}

# ============================================================================
# ФУНКЦИИ ДЛЯ СОЗДАНИЯ ЭЛЕМЕНТОВ
# ============================================================================

def create_presentation():
    """Создаёт презентацию с одним слайдом"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_title(slide, text):
    """Добавляет заголовок слайда"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = text
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary_blue']
    title_para.font.name = 'Arial'

def add_column_box(slide, left, top, width, height, content):
    """Добавляет колонку с контентом"""
    # Основной блок
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['bg_light']
    shape.line.color.rgb = COLORS['primary_blue']
    shape.line.width = Pt(2)
    
    # Текст внутри блока
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.3)
    text_frame.margin_bottom = Inches(0.3)
    
    # Заголовок с иконкой
    title_para = text_frame.paragraphs[0]
    title_para.text = f"{content['icon']} {content['title']}"
    title_para.font.size = Pt(18)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['accent_pink']
    title_para.font.name = 'Arial'
    title_para.space_after = Pt(12)
    
    # Пункты списка
    for item in content['items']:
        para = text_frame.add_paragraph()
        para.text = item
        para.font.size = Pt(14)
        para.font.color.rgb = COLORS['text_dark']
        para.font.name = 'Arial'
        para.space_after = Pt(8)
        para.line_spacing = 1.3

def add_footer(slide, text):
    """Добавляет нижний блок с выводом"""
    footer_top = Inches(5.8)
    footer_height = Inches(1.4)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), footer_top, Inches(12.333), footer_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary_blue']
    shape.line.color.rgb = COLORS['primary_blue']
    shape.line.width = Pt(0)
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.5)
    text_frame.margin_right = Inches(0.5)
    text_frame.margin_top = Inches(0.3)
    text_frame.margin_bottom = Inches(0.3)
    
    para = text_frame.paragraphs[0]
    para.text = text
    para.font.size = Pt(16)
    para.font.bold = True
    para.font.color.rgb = COLORS['white']
    para.font.name = 'Arial'
    para.alignment = PP_ALIGN.CENTER

def add_tech_logos(slide):
    """Добавляет текстовые метки технологий (вместо иконок)"""
    logos = ['Python', 'Django', 'PostgreSQL', 'REST API', 'QR']
    x_start = Inches(0.5)
    y_pos = Inches(5.3)
    
    for i, logo in enumerate(logos):
        left = x_start + Inches(i * 2.4)
        textbox = slide.shapes.add_textbox(left, y_pos, Inches(2.2), Inches(0.4))
        frame = textbox.text_frame
        para = frame.paragraphs[0]
        para.text = f"⚡ {logo}"
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.color.rgb = COLORS['accent_pink']
        para.font.name = 'Arial'
        para.alignment = PP_ALIGN.CENTER

def add_github_qr_placeholder(slide):
    """Добавляет заглушку для QR-кода GitHub"""
    # ✅ ИСПРАВЛЕНО: используем RECTANGLE вместо SQUARE
    qr_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(11.5), Inches(0.3), Inches(1.3), Inches(1.3)
    )
    qr_box.fill.solid()
    qr_box.fill.fore_color.rgb = COLORS['white']
    qr_box.line.color.rgb = COLORS['primary_blue']
    qr_box.line.width = Pt(2)
    
    text_frame = qr_box.text_frame
    text_frame.word_wrap = True
    para = text_frame.paragraphs[0]
    para.text = "QR\nGitHub"
    para.font.size = Pt(10)
    para.font.color.rgb = COLORS['primary_blue']
    para.font.name = 'Arial'
    para.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# ============================================================================
# ГЛАВНАЯ ФУНКЦИЯ
# ============================================================================

def main():
    print("🚀 Создание слайда VetMis Technology...")
    
    # Создаём презентацию
    prs = create_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Добавляем элементы
    add_title(slide, SLIDE_CONTENT['title'])
    
    # Три колонки
    col_width = Inches(3.9)
    col_height = Inches(4.8)
    top_pos = Inches(1.4)
    
    add_column_box(slide, Inches(0.5), top_pos, col_width, col_height, SLIDE_CONTENT['column1'])
    add_column_box(slide, Inches(4.7), top_pos, col_width, col_height, SLIDE_CONTENT['column2'])
    add_column_box(slide, Inches(8.9), top_pos, col_width, col_height, SLIDE_CONTENT['column3'])
    
    # Логотипы технологий
    add_tech_logos(slide)
    
    # Заглушка QR-кода
    add_github_qr_placeholder(slide)
    
    # Нижний блок с выводом
    add_footer(slide, SLIDE_CONTENT['footer'])
    
    # Сохраняем файл
    filename = 'VetMis_Technology_Slide.pptx'
    prs.save(filename)
    
    print(f"✅ Слайд успешно создан: {filename}")
    print("\n📥 Инструкция:")
    print("1. Откройте файл в PowerPoint")
    print("2. Скопируйте слайд (Ctrl+C)")
    print("3. Вставьте в вашу презентацию (Ctrl+V)")
    print("\n🔗 GitHub репозиторий: https://github.com/danya257/vetmis2")

if __name__ == '__main__':
    main()